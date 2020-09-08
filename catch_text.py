# -*- coding: utf-8 -*-
# !python3
"""
PPT catch text
"""

from pptx import Presentation

def ppt_catch_text(filename):
    """
    Extract all text from slides in presentation
    """
    prs = Presentation(filename)
    text = {}

    for x in range(len(prs.slides)):
        text[x] = []
        for shape in prs.slides[x].shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text[x].append(run.text)

    return text

print(ppt_catch_text('./sample.pptx'))
