# -*- coding: utf-8 -*-
# !python3
"""
PPT catch format text
"""

from pptx import Presentation

from pptx.enum.shapes import MSO_SHAPE_TYPE

def ppt_catch_format_text(filename):
    """
    Extract all text from the slides in the presentation 
    and return it in the format.
    """
    prs = Presentation(filename)
    txt_oa = {}
    for x in range(len(prs.slides)):
        txt_oa[x] = []

        # Only on table elements
        for shape in prs.slides[x].shapes:
            if hasattr(shape, "table"):
                for row in shape.table.rows:
                    row_str = ""
                    for cell in row.cells:
                        row_str += cell.text_frame.text + " | "
                    row_text = row_str.encode('utf-8').strip().decode()
                    txt_oa[x].append(row_text)

        # Only on text-boxes outside group elements
        for shape in prs.slides[x].shapes:
            if hasattr(shape, "text"):
                row_text_arr = shape.text.encode('utf-8').strip().decode().split("\n")
                for row_text in row_text_arr:
                    txt_oa[x].append(row_text)

        # Only on group shapes elements
        group_shapes = [shp for shp in prs.slides[x].shapes 
                        if shp.shape_type ==MSO_SHAPE_TYPE.GROUP]
        for group_shape in group_shapes:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    row_text = shape.text.encode('utf-8').strip().decode()
                    txt_oa[x].append(row_text)
    return txt_oa

print(ppt_catch_format_text('./sample.pptx'))
